"""Office document preview provider (.docx, .xlsx, .pptx)."""
from __future__ import annotations

import html
from pathlib import Path

from biome_fm.preview.provider import ContentKind, PreviewRequest, PreviewResult

_MAX_BYTES = 2 * 1024 * 1024
_EXTS = {".docx", ".xlsx", ".pptx"}


class OfficeProvider:
    priority = 3

    def can_handle(self, path: Path) -> bool:
        return path.suffix.lower() in _EXTS

    def render(self, req: PreviewRequest) -> PreviewResult:
        path = req.path
        try:
            if path.stat().st_size > _MAX_BYTES:
                return PreviewResult(ContentKind.ERROR, "File too large for preview")
            ext = path.suffix.lower()
            if ext == ".docx":
                return self._docx(path)
            if ext == ".xlsx":
                return self._xlsx(path)
            return self._pptx(path)
        except ImportError as e:
            return PreviewResult(ContentKind.ERROR, f"Missing dependency: {e}")
        except Exception as e:
            return PreviewResult(ContentKind.ERROR, f"Preview failed: {e}")

    def _docx(self, path: Path) -> PreviewResult:
        from docx import Document  # type: ignore[import]
        doc = Document(str(path))
        lines = [html.escape(p.text) for p in doc.paragraphs if p.text.strip()]
        return PreviewResult(ContentKind.HTML, "<br>".join(lines) or "Empty document")

    def _xlsx(self, path: Path) -> PreviewResult:
        from openpyxl import load_workbook  # type: ignore[import]
        wb = load_workbook(str(path), read_only=True)
        ws = wb.active
        rows = list(ws.iter_rows(max_row=50, values_only=True))
        cells = lambda row: "".join(f"<td>{html.escape(str(v or ''))}</td>" for v in row)
        tbody = "".join(f"<tr>{cells(r)}</tr>" for r in rows)
        return PreviewResult(ContentKind.HTML, f"<table>{tbody}</table>" or "Empty sheet")

    def _pptx(self, path: Path) -> PreviewResult:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                html.escape(tf.text)
                for shape in slide.shapes
                if shape.has_text_frame
                for tf in [shape.text_frame]
                if tf.text.strip()
            ]
            if texts:
                parts.append(f"<b>Slide {i}</b><br>" + "<br>".join(texts))
        return PreviewResult(ContentKind.HTML, "<hr>".join(parts) or "Empty presentation")
